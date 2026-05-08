from pathlib import Path
import textwrap

from PIL import Image, ImageEnhance, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "diploma_output"
SCREEN_DIR = OUTPUT_DIR / "screenshots"
ASSET_DIR = OUTPUT_DIR / "presentation_assets"
PPTX_PATH = OUTPUT_DIR / "Презентация_Платформа_корпоративного_обучения.pptx"
SPEECH_PATH = OUTPUT_DIR / "Краткий_рассказ_к_презентации.txt"


COLORS = {
    "ink": "10211B",
    "muted": "5D6F68",
    "subtle": "7A8E88",
    "brand": "0F7A65",
    "brand_dark": "0B5A4B",
    "accent": "D98930",
    "accent_soft": "FFE8C2",
    "cream": "FBF6EF",
    "surface": "FFFFFF",
    "surface_2": "F4FAF7",
    "line": "D5DDD7",
    "dark": "0E1F1A",
}


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def emu_to_inches(value):
    return value / 914400


def add_rect(slide, x, y, w, h, fill, line=None, radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color="ink",
    bold=False,
    font="Segoe UI",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.08,
    auto_fit=False,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if auto_fit else MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(COLORS[color])
    return box


def add_bullets(slide, items, x, y, w, h, size=17, color="ink"):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COLORS[color])
        p.line_spacing = 1.08
        p.space_after = Pt(5)
        p._p.get_or_add_pPr().set("marL", "342900")
        p._p.get_or_add_pPr().set("hanging", "171450")
    return box


def add_footer(slide, prs, number, title="Платформа корпоративного обучения"):
    y = prs.slide_height - Inches(0.38)
    add_text(slide, title, Inches(0.55), y, Inches(5.2), Inches(0.18), size=8.5, color="subtle")
    add_text(
        slide,
        f"{number:02d}",
        prs.slide_width - Inches(0.83),
        y,
        Inches(0.3),
        Inches(0.18),
        size=8.5,
        color="subtle",
        align=PP_ALIGN.RIGHT,
    )


def add_title(slide, eyebrow, title, subtitle=None):
    add_text(slide, eyebrow.upper(), Inches(0.65), Inches(0.45), Inches(8.5), Inches(0.28), size=10.5, color="brand", bold=True)
    add_text(slide, title, Inches(0.65), Inches(0.82), Inches(9.1), Inches(0.82), size=31, color="ink", bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.67), Inches(1.63), Inches(9.3), Inches(0.5), size=15, color="muted")


def add_chip(slide, text, x, y, w=None, fill="surface_2", color="brand", line="line"):
    width = w or Inches(max(1.0, min(2.6, 0.15 * len(text) + 0.45)))
    add_rect(slide, x, y, width, Inches(0.35), COLORS[fill], COLORS[line])
    add_text(slide, text, x + Inches(0.11), y + Inches(0.075), width - Inches(0.22), Inches(0.13), size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    return width


def add_card(slide, x, y, w, h, title, body, accent="brand", number=None):
    add_rect(slide, x, y, w, h, COLORS["surface"], COLORS["line"])
    add_rect(slide, x, y, Inches(0.08), h, COLORS[accent], None, MSO_SHAPE.RECTANGLE)
    if number:
        add_text(slide, number, x + Inches(0.22), y + Inches(0.22), Inches(0.45), Inches(0.3), size=15, color=accent, bold=True)
        title_x = x + Inches(0.75)
        title_w = w - Inches(0.95)
    else:
        title_x = x + Inches(0.25)
        title_w = w - Inches(0.5)
    add_text(slide, title, title_x, y + Inches(0.19), title_w, Inches(0.36), size=15, color="ink", bold=True)
    add_text(slide, body, x + Inches(0.25), y + Inches(0.66), w - Inches(0.5), h - Inches(0.82), size=11.5, color="muted", line_spacing=1.05)


def add_image(slide, path, x, y, w, h, border=True):
    if border:
        add_rect(slide, x - Inches(0.03), y - Inches(0.03), w + Inches(0.06), h + Inches(0.06), COLORS["surface"], COLORS["line"])
    slide.shapes.add_picture(str(path), x, y, w, h)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text.strip()


def cover_crop(src: Path, dst: Path, ratio: float, anchor_y=0.0, anchor_x=0.5, enhance=True):
    image = Image.open(src).convert("RGB")
    width, height = image.size
    current_ratio = width / height

    if current_ratio > ratio:
        new_width = int(height * ratio)
        left = int((width - new_width) * anchor_x)
        box = (left, 0, left + new_width, height)
    else:
        new_height = int(width / ratio)
        top = int((height - new_height) * anchor_y)
        box = (0, top, width, top + new_height)

    cropped = image.crop(box)
    cropped = ImageOps.expand(cropped, border=2, fill="#d5ddd7")
    if enhance:
        cropped = ImageEnhance.Sharpness(cropped).enhance(1.08)
        cropped = ImageEnhance.Contrast(cropped).enhance(1.02)
    cropped.save(dst, quality=94)
    return dst


def prepare_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    mapping = {
        "login": ("01-login.png", 16 / 9, 0.0),
        "dashboard": ("02-dashboard.png", 16 / 9, 0.0),
        "users": ("03-users.png", 16 / 10, 0.0),
        "chat": ("04-chat.png", 16 / 10, 0.0),
        "roles": ("05-roles.png", 16 / 9, 0.0),
        "course": ("06-course-viewer.png", 16 / 10, 0.0),
        "courses": ("07-courses.png", 16 / 9, 0.0),
        "certificates": ("08-certificates.png", 16 / 10, 0.0),
    }
    result = {}
    for key, (name, ratio, anchor_y) in mapping.items():
        result[key] = cover_crop(SCREEN_DIR / name, ASSET_DIR / f"{key}.jpg", ratio, anchor_y=anchor_y)
    return result


def base_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, COLORS["cream"], None, MSO_SHAPE.RECTANGLE)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.14), COLORS["brand"], None, MSO_SHAPE.RECTANGLE)
    add_rect(slide, 0, prs.slide_height - Inches(0.08), prs.slide_width, Inches(0.08), COLORS["accent"], None, MSO_SHAPE.RECTANGLE)
    return slide


def connector(slide, x1, y1, x2, y2, color="line", width=1.4):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = rgb(COLORS[color])
    line.line.width = Pt(width)
    return line


def build_presentation():
    assets = prepare_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Платформа для корпоративного обучения"
    prs.core_properties.subject = "Дипломная презентация"
    prs.core_properties.author = "Codex"

    # 1
    slide = base_slide(prs)
    add_text(slide, "ДИПЛОМНЫЙ ПРОЕКТ", Inches(0.72), Inches(0.62), Inches(4.0), Inches(0.28), size=10.5, color="brand", bold=True)
    add_text(
        slide,
        "Платформа для\nкорпоративного\nобучения",
        Inches(0.70),
        Inches(1.07),
        Inches(5.2),
        Inches(2.05),
        size=37,
        color="ink",
        bold=True,
        line_spacing=0.92,
    )
    add_text(
        slide,
        "Многоарендная LMS для управления курсами, сотрудниками, прогрессом, коммуникациями, отчетами и сертификатами.",
        Inches(0.72),
        Inches(3.30),
        Inches(5.05),
        Inches(0.85),
        size=16,
        color="muted",
    )
    x = Inches(0.72)
    for chip in ["React", "Node.js", "Express", "Prisma", "SQLite", "JWT", "Socket.IO"]:
        w = add_chip(slide, chip, x, Inches(4.48), fill="surface")
        x += w + Inches(0.11)
    add_rect(slide, Inches(6.28), Inches(0.75), Inches(6.28), Inches(5.46), COLORS["surface"], COLORS["line"])
    add_image(slide, assets["dashboard"], Inches(6.45), Inches(0.92), Inches(5.94), Inches(3.34), border=False)
    add_text(slide, "Рабочее пространство LMS", Inches(6.55), Inches(4.55), Inches(5.6), Inches(0.35), size=19, color="ink", bold=True)
    add_text(slide, "Единый экран для обучения, чата, статуса курсов, сертификатов и управленческой сводки.", Inches(6.55), Inches(4.95), Inches(5.55), Inches(0.65), size=12.5, color="muted")
    add_footer(slide, prs, 1)
    add_notes(
        slide,
        "Добрый день. Я представляю дипломный проект: платформу для корпоративного обучения. Это web-приложение формата LMS, которое помогает компании централизованно управлять курсами, сотрудниками, прогрессом обучения, коммуникациями, отчетами и сертификатами.",
    )

    # 2
    slide = base_slide(prs)
    add_title(slide, "Актуальность", "Исходная проблема", "Обучение в небольшой IT-компании часто остается ручным процессом, хотя напрямую влияет на адаптацию и качество работы.")
    problems = [
        ("Материалы разрознены", "Учебные файлы, ссылки и регламенты хранятся в разных каналах."),
        ("Контроль вручную", "Руководителям сложно быстро понять, кто прошел обязательные курсы."),
        ("Нет единой коммуникации", "Вопросы по обучению теряются в мессенджерах и почте."),
        ("Слабая отчетность", "Сертификаты и управленческие отчеты формируются с лишними трудозатратами."),
    ]
    positions = [(0.72, 2.42), (4.07, 2.42), (0.72, 4.25), (4.07, 4.25)]
    for idx, ((title, body), (x, y)) in enumerate(zip(problems, positions), 1):
        add_card(slide, Inches(x), Inches(y), Inches(3.02), Inches(1.35), title, body, accent="accent" if idx % 2 == 0 else "brand", number=f"{idx}")
    add_rect(slide, Inches(7.65), Inches(2.35), Inches(4.9), Inches(3.42), COLORS["surface"], COLORS["line"])
    add_text(slide, "AS-IS", Inches(7.95), Inches(2.65), Inches(1.0), Inches(0.25), size=12, color="accent", bold=True)
    flow = [
        ("Учебные материалы", 8.0, 3.15),
        ("Почта и мессенджеры", 10.05, 3.15),
        ("Ручные таблицы", 8.0, 4.45),
        ("Неочевидный прогресс", 10.05, 4.45),
    ]
    for label, x0, y0 in flow:
        add_rect(slide, Inches(x0), Inches(y0), Inches(1.65), Inches(0.67), COLORS["surface_2"], COLORS["line"])
        add_text(slide, label, Inches(x0 + 0.1), Inches(y0 + 0.17), Inches(1.45), Inches(0.22), size=9.4, color="ink", bold=True, align=PP_ALIGN.CENTER)
    connector(slide, Inches(9.65), Inches(3.49), Inches(10.05), Inches(3.49), "accent")
    connector(slide, Inches(8.82), Inches(3.82), Inches(8.82), Inches(4.45), "accent")
    connector(slide, Inches(9.65), Inches(4.79), Inches(10.05), Inches(4.79), "accent")
    add_text(slide, "Итог: процесс зависит от людей и каналов связи, а не от единого информационного контура.", Inches(7.95), Inches(5.35), Inches(4.05), Inches(0.44), size=11.5, color="muted")
    add_footer(slide, prs, 2)
    add_notes(
        slide,
        "Проблема проекта в том, что обучение сотрудников обычно распределено между файлами, переписками и ручными таблицами. Из-за этого сложно контролировать прохождение курсов, подтверждать результат и быстро готовить отчеты. Поэтому появляется потребность в единой платформе.",
    )

    # 3
    slide = base_slide(prs)
    add_title(slide, "Цель и задачи", "Что нужно было разработать", "Цель проекта — создать локально разворачиваемую LMS-платформу для автоматизации корпоративного обучения.")
    add_rect(slide, Inches(0.72), Inches(2.33), Inches(4.25), Inches(3.68), COLORS["dark"], None)
    add_text(slide, "Цель", Inches(1.05), Inches(2.70), Inches(0.9), Inches(0.28), size=12, color="accent", bold=True)
    add_text(
        slide,
        "Разработать web-платформу, которая объединяет создание курсов, назначение обучения, фиксацию прогресса, коммуникацию, сертификаты и отчеты в одном защищенном контуре.",
        Inches(1.05),
        Inches(3.22),
        Inches(3.45),
        Inches(1.85),
        size=17,
        color="surface",
        bold=True,
        line_spacing=1.03,
    )
    tasks = [
        ("1", "Исследовать процесс обучения и ограничения текущего состояния"),
        ("2", "Определить роли, права доступа и требования к данным"),
        ("3", "Спроектировать архитектуру клиента, API и базы данных"),
        ("4", "Реализовать курсы, прогресс, чат, отчеты и сертификаты"),
        ("5", "Подготовить систему к локальному запуску и дальнейшему развитию"),
    ]
    y = Inches(2.24)
    for num, text in tasks:
        add_rect(slide, Inches(5.45), y, Inches(6.58), Inches(0.62), COLORS["surface"], COLORS["line"])
        add_text(slide, num, Inches(5.68), y + Inches(0.15), Inches(0.32), Inches(0.18), size=11.5, color="brand", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, text, Inches(6.15), y + Inches(0.13), Inches(5.55), Inches(0.22), size=12.3, color="ink", bold=True)
        y += Inches(0.74)
    add_footer(slide, prs, 3)
    add_notes(
        slide,
        "Цель работы состоит не только в создании каталога учебных материалов. Система должна закрывать полный цикл: от регистрации компании и настройки ролей до назначения курсов, прохождения обучения, коммуникаций, выдачи сертификата и формирования управленческой отчетности.",
    )

    # 4
    slide = base_slide(prs)
    add_title(slide, "Концепция", "Пользователи и права доступа", "Платформа построена вокруг компании-арендатора и ролевой модели RBAC.")
    add_rect(slide, Inches(4.95), Inches(2.22), Inches(3.15), Inches(1.55), COLORS["dark"], None)
    add_text(slide, "Компания", Inches(5.25), Inches(2.58), Inches(2.55), Inches(0.32), size=21, color="surface", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "companyId изолирует пользователей,\nкурсы, роли, чаты и отчеты", Inches(5.18), Inches(3.08), Inches(2.7), Inches(0.35), size=10.5, color="accent_soft", align=PP_ALIGN.CENTER)
    roles = [
        ("Администратор", "компания, пользователи, роли,\nнастройки, отчеты", 0.85, 2.15, "brand"),
        ("HR", "сотрудники, отделы,\nназначения, прогресс", 0.85, 4.35, "accent"),
        ("Менеджер", "курсы команды,\nчат, контроль сроков", 9.05, 2.15, "brand"),
        ("Сотрудник", "личные курсы,\nвопросы, сертификаты", 9.05, 4.35, "accent"),
    ]
    for title, body, x, y, accent in roles:
        add_card(slide, Inches(x), Inches(y), Inches(3.0), Inches(1.35), title, body, accent=accent)
    connector(slide, Inches(4.95), Inches(2.98), Inches(3.85), Inches(2.82), "line")
    connector(slide, Inches(4.95), Inches(3.35), Inches(3.85), Inches(5.02), "line")
    connector(slide, Inches(8.10), Inches(2.98), Inches(9.05), Inches(2.82), "line")
    connector(slide, Inches(8.10), Inches(3.35), Inches(9.05), Inches(5.02), "line")
    add_rect(slide, Inches(4.05), Inches(4.78), Inches(4.95), Inches(0.84), COLORS["surface"], COLORS["line"])
    add_text(slide, "RBAC: права хранятся в JSON-представлении и проверяются на сервере и клиенте", Inches(4.35), Inches(5.04), Inches(4.35), Inches(0.2), size=11.2, color="ink", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, 4)
    add_notes(
        slide,
        "Важная особенность решения — многоарендная модель. Одна платформа может обслуживать несколько компаний, но данные каждой компании логически изолированы через companyId. Доступ внутри компании регулируется ролями: администратор, HR, менеджер и сотрудник.",
    )

    # 5
    slide = base_slide(prs)
    add_title(slide, "Функциональность", "Что реализовано в системе", "MVP закрывает основной цикл корпоративного обучения и сопровождающие административные процессы.")
    modules = [
        ("Регистрация компании", "создание арендатора, базовых ролей и первого администратора"),
        ("Пользователи и отделы", "профили, статусы, роли, принадлежность к подразделениям"),
        ("Курсы и уроки", "модули, материалы, видео, файлы, тесты и статусы курсов"),
        ("Назначения и прогресс", "дедлайны, завершение уроков, проценты прохождения"),
        ("Чат", "диалоги сотрудников и руководителей, вложения, realtime-события"),
        ("Сертификаты", "HTML/PDF, код проверки, подпись руководителя"),
        ("Дашборд и отчеты", "сводка по обучению, CSV/JSON/PDF-выгрузки"),
        ("Аудит и задачи", "журнал действий, фоновые операции, файловое хранилище"),
    ]
    for i, (title, body) in enumerate(modules):
        row = i // 4
        col = i % 4
        x = Inches(0.63 + col * 3.12)
        y = Inches(2.20 + row * 1.72)
        add_card(slide, x, y, Inches(2.72), Inches(1.25), title, body, accent="brand" if i % 2 == 0 else "accent")
    add_rect(slide, Inches(0.72), Inches(5.92), Inches(11.88), Inches(0.58), COLORS["surface_2"], COLORS["line"])
    add_text(slide, "Сценарий: создать компанию → настроить роли → добавить сотрудников → создать курс → назначить обучение → пройти уроки → получить отчет и сертификат", Inches(1.0), Inches(6.11), Inches(11.3), Inches(0.2), size=11.5, color="ink", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, 5)
    add_notes(
        slide,
        "В реализованной версии есть регистрация компании, управление пользователями, ролями и отделами, создание курсов с модулями и уроками, назначение обучения, учет прогресса, чат, отчеты, аудит и генерация сертификатов. То есть система закрывает не отдельную функцию, а рабочий процесс целиком.",
    )

    # 6
    slide = base_slide(prs)
    add_title(slide, "Архитектура", "Клиент-серверная реализация", "Решение разделено на React-клиент, Express API, сервисный слой и Prisma-модель данных.")
    tiers = [
        ("Клиент", "React + Vite + Tailwind\nмаршрутизация, формы, защищенные экраны", 0.85, 2.35, "brand"),
        ("API", "Node.js + Express\nroutes → middleware → controllers", 4.05, 2.35, "accent"),
        ("Бизнес-логика", "services\nпользователи, курсы, чат, отчеты, сертификаты", 7.25, 2.35, "brand"),
        ("Данные и файлы", "Prisma + SQLite + uploads\nмодель, сертификаты, материалы, вложения", 10.45, 2.35, "accent"),
    ]
    for title, body, x, y, accent in tiers:
        add_rect(slide, Inches(x), Inches(y), Inches(2.42), Inches(1.42), COLORS["surface"], COLORS["line"])
        add_rect(slide, Inches(x), Inches(y), Inches(2.42), Inches(0.17), COLORS[accent], None, MSO_SHAPE.RECTANGLE)
        add_text(slide, title, Inches(x + 0.18), Inches(y + 0.35), Inches(2.05), Inches(0.28), size=16, color="ink", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, Inches(x + 0.18), Inches(y + 0.80), Inches(2.05), Inches(0.38), size=9.8, color="muted", align=PP_ALIGN.CENTER)
    for x in [3.28, 6.48, 9.68]:
        connector(slide, Inches(x), Inches(3.06), Inches(x + 0.55), Inches(3.06), "brand", width=1.8)
    add_rect(slide, Inches(1.1), Inches(4.85), Inches(5.35), Inches(0.88), COLORS["dark"], None)
    add_text(slide, "MVC в проекте", Inches(1.45), Inches(5.08), Inches(1.65), Inches(0.25), size=13, color="accent", bold=True)
    add_text(slide, "View: React-страницы • Controller: Express routes/controllers • Model: Prisma + services", Inches(3.05), Inches(5.08), Inches(3.0), Inches(0.28), size=10.2, color="surface")
    add_rect(slide, Inches(7.0), Inches(4.85), Inches(5.0), Inches(0.88), COLORS["surface"], COLORS["line"])
    x = Inches(7.23)
    for chip in ["JWT", "bcrypt", "Zod", "Multer", "PDFKit", "Socket.IO"]:
        w = add_chip(slide, chip, x, Inches(5.10), fill="surface_2", color="brand")
        x += w + Inches(0.09)
    add_footer(slide, prs, 6)
    add_notes(
        slide,
        "Технически система построена по клиент-серверной архитектуре. На клиенте используется React-приложение, серверная часть реализована на Node.js и Express. Логика вынесена в сервисы, а работа с базой организована через Prisma поверх SQLite. Для файлов используется локальный каталог uploads, для чата — Socket.IO.",
    )

    # 7
    slide = base_slide(prs)
    add_title(slide, "Данные", "Модель предметной области", "Все ключевые сущности связаны с компанией и не смешиваются между арендаторами.")
    add_rect(slide, Inches(5.28), Inches(2.05), Inches(2.75), Inches(0.85), COLORS["dark"], None)
    add_text(slide, "Company", Inches(5.65), Inches(2.32), Inches(2.0), Inches(0.22), size=18, color="surface", bold=True, align=PP_ALIGN.CENTER)
    groups = [
        ("Оргструктура", "Users\nRoles\nDepartments", 0.9, 3.05, "brand"),
        ("Обучение", "Courses\nModules\nLessons\nQuizzes", 3.4, 4.35, "accent"),
        ("Исполнение", "Assignments\nProgress\nCertificates", 6.1, 4.35, "brand"),
        ("Коммуникации", "ChatThreads\nParticipants\nMessages", 8.85, 4.35, "accent"),
        ("Сопровождение", "AuditLogs\nBackgroundJobs\nUploads", 10.35, 3.05, "brand"),
    ]
    for title, body, x, y, accent in groups:
        add_rect(slide, Inches(x), Inches(y), Inches(2.15), Inches(1.12), COLORS["surface"], COLORS["line"])
        add_text(slide, title, Inches(x + 0.16), Inches(y + 0.18), Inches(1.83), Inches(0.22), size=12.4, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, Inches(x + 0.16), Inches(y + 0.50), Inches(1.83), Inches(0.38), size=10.5, color="ink", align=PP_ALIGN.CENTER, line_spacing=1.0)
        connector(slide, Inches(6.65), Inches(2.90), Inches(x + 1.08), Inches(y), "line")
    add_rect(slide, Inches(0.9), Inches(5.92), Inches(11.55), Inches(0.52), COLORS["surface_2"], COLORS["line"])
    add_text(slide, "Уникальные ограничения и индексы учитывают companyId: email, роли, отделы, курсы, отчеты и права доступны только внутри своей компании.", Inches(1.15), Inches(6.08), Inches(11.0), Inches(0.2), size=11.1, color="ink", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, 7)
    add_notes(
        slide,
        "Модель данных строится вокруг сущности Company. С ней связаны пользователи, роли, отделы, курсы, уроки, тесты, назначения, прогресс, сообщения, сертификаты, аудит и фоновые задачи. Такая структура позволяет масштабировать систему на несколько компаний и сохранять логическую изоляцию данных.",
    )

    # 8
    slide = base_slide(prs)
    add_title(slide, "Интерфейс", "Рабочее пространство пользователя", "Дашборд объединяет обучение, коммуникации и быстрые действия.")
    add_image(slide, assets["dashboard"], Inches(0.75), Inches(2.22), Inches(7.15), Inches(4.02))
    add_card(slide, Inches(8.35), Inches(2.25), Inches(3.9), Inches(1.0), "Сводка обучения", "назначенные, активные и завершенные курсы видны на первом экране", accent="brand")
    add_card(slide, Inches(8.35), Inches(3.55), Inches(3.9), Inches(1.0), "Коммуникации", "чат и новые обращения обновляют состояние интерфейса", accent="accent")
    add_card(slide, Inches(8.35), Inches(4.85), Inches(3.9), Inches(1.0), "Результаты", "сертификаты, отчеты и действия доступны из единого кабинета", accent="brand")
    add_footer(slide, prs, 8)
    add_notes(
        slide,
        "Интерфейс сделан как единое рабочее пространство. На дашборде пользователь видит активные курсы, сообщения, сертификаты и важные показатели. Для сотрудника это личный кабинет обучения, для руководителя — инструмент контроля команды.",
    )

    # 9
    slide = base_slide(prs)
    add_title(slide, "Администрирование", "Управление сотрудниками, ролями и курсами", "Администратор и HR могут поддерживать организационную структуру и учебные программы без обращения к базе данных.")
    add_image(slide, assets["users"], Inches(0.74), Inches(2.15), Inches(5.2), Inches(3.25))
    add_image(slide, assets["courses"], Inches(6.26), Inches(2.15), Inches(5.78), Inches(3.25))
    add_rect(slide, Inches(0.95), Inches(5.66), Inches(10.92), Inches(0.58), COLORS["surface"], COLORS["line"])
    add_text(slide, "Карточка сотрудника связывает профиль, роль, отдел, статус доступа, назначенное обучение и историю завершенных курсов.", Inches(1.2), Inches(5.85), Inches(10.45), Inches(0.2), size=11.3, color="ink", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, 9)
    add_notes(
        slide,
        "Административная часть позволяет управлять сотрудниками, ролями, отделами и курсами. Это важно, потому что LMS должна быть не только экраном для прохождения уроков, но и инструментом сопровождения организационной структуры.",
    )

    # 10
    slide = base_slide(prs)
    add_title(slide, "Пользовательский цикл", "Обучение, чат и сертификаты", "Система ведет сотрудника от назначения курса до подтверждения результата.")
    add_image(slide, assets["course"], Inches(0.72), Inches(2.15), Inches(3.72), Inches(2.33))
    add_image(slide, assets["chat"], Inches(4.83), Inches(2.15), Inches(3.72), Inches(2.33))
    add_image(slide, assets["certificates"], Inches(8.94), Inches(2.15), Inches(3.72), Inches(2.33))
    captions = [
        ("1. Прохождение курса", "модули, уроки, файлы и отметка прогресса"),
        ("2. Вопросы в чате", "диалоги с руководителем и realtime-обновления"),
        ("3. Сертификат", "PDF, код проверки и подпись руководителя"),
    ]
    for i, (title, body) in enumerate(captions):
        x = Inches(0.85 + i * 4.11)
        add_text(slide, title, x, Inches(4.82), Inches(3.45), Inches(0.25), size=14.5, color="ink", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x, Inches(5.18), Inches(3.45), Inches(0.36), size=10.5, color="muted", align=PP_ALIGN.CENTER)
    add_footer(slide, prs, 10)
    add_notes(
        slide,
        "Основной пользовательский цикл выглядит так: сотруднику назначают курс, он проходит уроки и фиксирует прогресс. Если возникают вопросы, он может обратиться в чат. После успешного завершения формируется сертификат с проверочным кодом и возможностью выгрузки PDF.",
    )

    # 11
    slide = base_slide(prs)
    add_title(slide, "Результат", "Что получено и как развивать дальше", "Реализована работоспособная версия LMS, пригодная для локального запуска и поэтапного расширения.")
    add_rect(slide, Inches(0.82), Inches(2.15), Inches(5.35), Inches(3.75), COLORS["surface"], COLORS["line"])
    add_text(slide, "Полученный результат", Inches(1.18), Inches(2.47), Inches(4.65), Inches(0.32), size=20, color="ink", bold=True)
    add_bullets(
        slide,
        [
            "многоарендная модель с изоляцией данных компаний",
            "RBAC, JWT-аутентификация и серверные проверки доступа",
            "управление курсами, уроками, назначениями и прогрессом",
            "чат, отчеты, аудит и PDF-сертификаты",
            "структура проекта готова к сопровождению и расширению",
        ],
        Inches(1.22),
        Inches(3.05),
        Inches(4.55),
        Inches(2.4),
        size=12.4,
    )
    add_rect(slide, Inches(7.05), Inches(2.15), Inches(5.35), Inches(3.75), COLORS["dark"], None)
    add_text(slide, "Дальнейшее развитие", Inches(7.40), Inches(2.47), Inches(4.65), Inches(0.32), size=20, color="surface", bold=True)
    add_bullets(
        slide,
        [
            "перенос базы данных на PostgreSQL",
            "замена локальных uploads на S3-совместимое хранилище",
            "расширение фоновых задач и уведомлений",
            "углубленная аналитика по обучению и компетенциям",
            "интеграция с корпоративной почтой или HR-системой",
        ],
        Inches(7.45),
        Inches(3.05),
        Inches(4.45),
        Inches(2.4),
        size=12.4,
        color="surface",
    )
    add_footer(slide, prs, 11)
    add_notes(
        slide,
        "Итогом работы стала рабочая версия платформы корпоративного обучения. Она решает задачу централизации учебного процесса и показывает основу для дальнейшего развития: переход на PostgreSQL, внешнее файловое хранилище, расширенные уведомления, аналитику и интеграции.",
    )

    # 12
    slide = base_slide(prs)
    add_title(slide, "Краткий рассказ", "Текст выступления", "Версию для чтения я также сохранил отдельным TXT-файлом рядом с презентацией.")
    speech_blocks = [
        ("1", "Проект посвящен разработке платформы корпоративного обучения для ООО «Эссеншиал Солюшнс»."),
        ("2", "Основная проблема — разрозненное хранение материалов, ручной контроль прогресса и отсутствие единого пространства для коммуникации и отчетности."),
        ("3", "Разработанная LMS объединяет роли, сотрудников, курсы, назначение обучения, прогресс, чат, отчеты и сертификаты."),
        ("4", "Архитектура построена на React, Node.js, Express, Prisma и SQLite; доступ защищен JWT и RBAC."),
        ("5", "Ключевая особенность — многоарендность: одна платформа может обслуживать несколько компаний без смешения данных."),
        ("6", "Результат — локально запускаемая версия системы, которую можно расширять до промышленной эксплуатации."),
    ]
    y = Inches(2.05)
    for num, text in speech_blocks:
        add_rect(slide, Inches(0.86), y, Inches(0.42), Inches(0.42), COLORS["brand"], None)
        add_text(slide, num, Inches(0.98), y + Inches(0.10), Inches(0.18), Inches(0.11), size=9, color="surface", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, text, Inches(1.52), y + Inches(0.06), Inches(10.6), Inches(0.2), size=13.0, color="ink", bold=True)
        y += Inches(0.65)
    add_footer(slide, prs, 12)
    add_notes(
        slide,
        "На этом слайде представлен сокращенный вариант рассказа. Полный связный текст выступления сохранен отдельным файлом в папке diploma_output и продублирован в заметках к презентации.",
    )

    prs.save(PPTX_PATH)
    return PPTX_PATH


SPEECH_TEXT = """\
Добрый день. Тема моего дипломного проекта — разработка платформы для корпоративного обучения на примере ООО «Эссеншиал Солюшнс».

Актуальность работы связана с тем, что во многих небольших IT-компаниях обучение сотрудников организовано через разрозненные материалы, переписки и ручной контроль. Из-за этого руководителям сложно быстро понять, какие курсы назначены, кто уже завершил обучение, где возникли вопросы и какие подтверждающие документы можно сформировать.

Цель проекта состояла в создании web-платформы, которая объединяет основные процессы корпоративного обучения: регистрацию компании, управление пользователями и ролями, создание курсов, назначение обучения, фиксацию прогресса, коммуникацию в чате, формирование отчетов и выдачу сертификатов.

Ключевая особенность решения — многоарендная модель. Одна технологическая платформа может обслуживать несколько компаний, но данные каждой компании изолируются через companyId. Внутри компании доступ регулируется ролевой моделью RBAC: администратор управляет настройками, HR и менеджеры контролируют сотрудников и обучение, а сотрудник проходит назначенные курсы и получает сертификаты.

Технически система реализована как клиент-серверное приложение. Клиентская часть построена на React, Vite и Tailwind CSS. Серверная часть использует Node.js и Express, бизнес-логика вынесена в сервисы, а модель данных описана через Prisma поверх SQLite. Для безопасности применяются JWT и bcrypt, для загрузки файлов используется Multer, для чата и обновлений — Socket.IO, для сертификатов и отчетов — генерация HTML/PDF.

В результате получена работоспособная версия LMS, пригодная для локального запуска и демонстрации полного цикла: администратор создает структуру компании и курс, назначает обучение сотруднику, сотрудник проходит уроки и задает вопросы в чате, а система фиксирует прогресс, формирует отчетность и выдает сертификат с кодом проверки.

Дальнейшее развитие проекта может включать перенос базы данных на PostgreSQL, подключение S3-совместимого файлового хранилища, расширение аналитики, уведомлений и интеграций с корпоративной почтой или HR-системами.
"""


def write_speech():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    SPEECH_PATH.write_text(textwrap.dedent(SPEECH_TEXT).strip() + "\n", encoding="utf-8")
    return SPEECH_PATH


if __name__ == "__main__":
    pptx_path = build_presentation()
    speech_path = write_speech()
    print(f"PowerPoint: {pptx_path}")
    print(f"Speech: {speech_path}")
